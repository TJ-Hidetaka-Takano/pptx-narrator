#!/usr/bin/env python3
"""PowerPointの発表者ノートをAivisSpeech Engineで音声化する。"""

from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import time
import wave
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import requests
from pptx import Presentation


@dataclass(frozen=True)
class VoiceSelection:
    style_id: int
    speaker_name: str
    style_name: str


def build_engine_session(base_url: str) -> requests.Session:
    """Avoid corporate proxies when talking to local/container-internal engine."""
    session = requests.Session()
    host = (urlparse(base_url).hostname or "").lower()
    if host in {"aivis-engine", "localhost", "127.0.0.1", "::1"}:
        session.trust_env = False
    return session


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="PPTXの発表者ノートをAivisSpeech EngineでWAV/MP3化します。"
    )
    parser.add_argument("pptx", nargs="?", type=Path, help="入力PPTXファイル")
    parser.add_argument(
        "--engine-url",
        default=os.getenv("AIVIS_ENGINE_URL", "http://aivis-engine:10101"),
        help="AivisSpeech EngineのURL",
    )
    parser.add_argument(
        "--output-dir", type=Path, default=Path("/work/output"), help="出力先"
    )
    parser.add_argument("--style-id", type=int, help="使用するスタイルID")
    parser.add_argument(
        "--list-voices", action="store_true", help="話者・スタイル一覧を表示して終了"
    )
    parser.add_argument(
        "--extract-only",
        action="store_true",
        help="発表者ノートを抽出し、音声生成は行わない",
    )
    parser.add_argument("--speed", type=float, default=1.0, help="話速倍率")
    parser.add_argument("--pitch", type=float, default=0.0, help="音高調整")
    parser.add_argument("--intonation", type=float, default=1.0, help="抑揚倍率")
    parser.add_argument(
        "--volume", type=float, default=1.0, help="音量倍率"
    )
    parser.add_argument(
        "--pre-phoneme-length", type=float, default=0.15, help="文頭無音秒"
    )
    parser.add_argument(
        "--post-phoneme-length", type=float, default=0.25, help="文末無音秒"
    )
    parser.add_argument(
        "--slide-gap", type=float, default=0.8, help="結合時のスライド間無音秒"
    )
    parser.add_argument(
        "--format", choices=("wav", "mp3"), default="mp3", help="個別音声の形式"
    )
    parser.add_argument(
        "--combine",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="全スライドを1本に結合する",
    )
    parser.add_argument(
        "--combined-name",
        default="lecture",
        help="結合ファイル名（拡張子を除く）",
    )
    parser.add_argument(
        "--include-slides",
        help="対象スライド。例: 1-5,8,10-12",
    )
    parser.add_argument(
        "--skip-empty", action=argparse.BooleanOptionalAction, default=True
    )
    parser.add_argument(
        "--wait-seconds",
        type=int,
        default=240,
        help="Engine起動待ちの最大秒数",
    )
    return parser.parse_args()


def wait_for_engine(base_url: str, timeout: int, session: requests.Session) -> None:
    deadline = time.monotonic() + timeout
    last_error = ""
    while time.monotonic() < deadline:
        try:
            response = session.get(f"{base_url.rstrip('/')}/version", timeout=5)
            response.raise_for_status()
            return
        except requests.RequestException as exc:
            last_error = str(exc)
            time.sleep(2)
    raise RuntimeError(f"AivisSpeech Engineへ接続できません: {last_error}")


def fetch_speakers(base_url: str, session: requests.Session) -> list[dict[str, Any]]:
    response = session.get(f"{base_url.rstrip('/')}/speakers", timeout=30)
    response.raise_for_status()
    data = response.json()
    if not isinstance(data, list):
        raise RuntimeError("/speakers の応答形式が不正です")
    return data


def list_voices(speakers: list[dict[str, Any]]) -> None:
    if not speakers:
        print("利用可能な音声モデルがありません。aivis-data/Models に .aivmx を配置してください。")
        return
    for speaker in speakers:
        speaker_name = speaker.get("name", "(名称なし)")
        for style in speaker.get("styles", []):
            print(
                f"style_id={style.get('id')}\t"
                f"speaker={speaker_name}\tstyle={style.get('name', '(名称なし)')}"
            )


def select_voice(
    speakers: list[dict[str, Any]], requested_style_id: int | None
) -> VoiceSelection:
    candidates: list[VoiceSelection] = []
    for speaker in speakers:
        for style in speaker.get("styles", []):
            candidates.append(
                VoiceSelection(
                    style_id=int(style["id"]),
                    speaker_name=str(speaker.get("name", "")),
                    style_name=str(style.get("name", "")),
                )
            )

    if not candidates:
        raise RuntimeError(
            "利用可能な音声モデルがありません。"
            "aivis-data/Models にライセンスを確認した .aivmx を配置してください。"
        )

    if requested_style_id is None:
        return candidates[0]

    for candidate in candidates:
        if candidate.style_id == requested_style_id:
            return candidate

    available = ", ".join(str(x.style_id) for x in candidates)
    raise RuntimeError(
        f"style_id={requested_style_id} は存在しません。利用可能: {available}"
    )


def parse_slide_spec(spec: str | None, total: int) -> set[int]:
    if not spec:
        return set(range(1, total + 1))

    result: set[int] = set()
    for part in spec.split(","):
        token = part.strip()
        if not token:
            continue
        if "-" in token:
            start_text, end_text = token.split("-", 1)
            start, end = int(start_text), int(end_text)
            if start > end:
                start, end = end, start
            result.update(range(start, end + 1))
        else:
            result.add(int(token))

    invalid = sorted(x for x in result if x < 1 or x > total)
    if invalid:
        raise ValueError(f"範囲外のスライド番号です: {invalid}")
    return result


def normalize_notes(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    text = "\n".join(line for line in lines if line)
    return text.strip()


def extract_notes(pptx_path: Path, include_spec: str | None) -> list[tuple[int, str]]:
    presentation = Presentation(str(pptx_path))
    included = parse_slide_spec(include_spec, len(presentation.slides))
    notes: list[tuple[int, str]] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide_number not in included or not slide.has_notes_slide:
            continue
        text_frame = slide.notes_slide.notes_text_frame
        if text_frame is None:
            continue
        notes.append((slide_number, normalize_notes(text_frame.text)))
    return notes


def write_notes_files(notes: list[tuple[int, str]], output_dir: Path) -> None:
    notes_dir = output_dir / "notes"
    notes_dir.mkdir(parents=True, exist_ok=True)
    for slide_number, text in notes:
        (notes_dir / f"slide-{slide_number:03d}.txt").write_text(
            text + ("\n" if text else ""), encoding="utf-8"
        )


def make_audio_query(
    base_url: str,
    text: str,
    style_id: int,
    args: argparse.Namespace,
    session: requests.Session,
) -> dict[str, Any]:
    response = session.post(
        f"{base_url.rstrip('/')}/audio_query",
        params={"text": text, "speaker": style_id},
        timeout=120,
    )
    response.raise_for_status()
    query = response.json()
    query["speedScale"] = args.speed
    query["pitchScale"] = args.pitch
    query["intonationScale"] = args.intonation
    query["volumeScale"] = args.volume
    query["prePhonemeLength"] = args.pre_phoneme_length
    query["postPhonemeLength"] = args.post_phoneme_length
    return query


def synthesize(
    base_url: str,
    query: dict[str, Any],
    style_id: int,
    output_wav: Path,
    session: requests.Session,
) -> None:
    response = session.post(
        f"{base_url.rstrip('/')}/synthesis",
        params={"speaker": style_id},
        json=query,
        timeout=600,
    )
    response.raise_for_status()
    output_wav.write_bytes(response.content)


def convert_audio(source_wav: Path, destination: Path) -> None:
    if destination.suffix.lower() == ".wav":
        shutil.copy2(source_wav, destination)
        return
    subprocess.run(
        [
            "ffmpeg",
            "-hide_banner",
            "-loglevel",
            "error",
            "-y",
            "-i",
            str(source_wav),
            "-codec:a",
            "libmp3lame",
            "-q:a",
            "2",
            str(destination),
        ],
        check=True,
    )


def wav_parameters(path: Path) -> tuple[int, int, int]:
    with wave.open(str(path), "rb") as wav_file:
        return (
            wav_file.getnchannels(),
            wav_file.getsampwidth(),
            wav_file.getframerate(),
        )


def create_silence(reference_wav: Path, duration: float, output_wav: Path) -> None:
    channels, sample_width, sample_rate = wav_parameters(reference_wav)
    frame_count = max(0, round(sample_rate * duration))
    with wave.open(str(output_wav), "wb") as wav_file:
        wav_file.setnchannels(channels)
        wav_file.setsampwidth(sample_width)
        wav_file.setframerate(sample_rate)
        wav_file.writeframes(b"\x00" * frame_count * channels * sample_width)


def combine_wavs(
    wav_files: list[Path],
    slide_gap: float,
    destination: Path,
    work_dir: Path,
) -> None:
    if not wav_files:
        return

    silence = work_dir / "slide-gap.wav"
    create_silence(wav_files[0], slide_gap, silence)

    concat_file = work_dir / "concat.txt"
    lines: list[str] = []
    for index, wav_file in enumerate(wav_files):
        lines.append(f"file '{wav_file.resolve().as_posix()}'")
        if index != len(wav_files) - 1 and slide_gap > 0:
            lines.append(f"file '{silence.resolve().as_posix()}'")
    concat_file.write_text("\n".join(lines) + "\n", encoding="utf-8")

    output_wav = work_dir / "combined.wav"
    subprocess.run(
        [
            "ffmpeg",
            "-hide_banner",
            "-loglevel",
            "error",
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(concat_file),
            "-c",
            "copy",
            str(output_wav),
        ],
        check=True,
    )
    convert_audio(output_wav, destination)


def main() -> int:
    args = parse_args()
    base_url = args.engine_url.rstrip("/")
    session = build_engine_session(base_url)

    if args.list_voices:
        wait_for_engine(base_url, args.wait_seconds, session)
        list_voices(fetch_speakers(base_url, session))
        return 0

    if args.pptx is None:
        print("エラー: PPTXファイルを指定してください。", file=sys.stderr)
        return 2
    if not args.pptx.is_file():
        print(f"エラー: ファイルが存在しません: {args.pptx}", file=sys.stderr)
        return 2

    args.output_dir.mkdir(parents=True, exist_ok=True)
    notes = extract_notes(args.pptx, args.include_slides)
    if args.skip_empty:
        notes = [(number, text) for number, text in notes if text]

    write_notes_files(notes, args.output_dir)
    print(f"発表者ノートを {len(notes)} スライドから抽出しました。")

    if args.extract_only:
        return 0

    wait_for_engine(base_url, args.wait_seconds, session)
    voice = select_voice(fetch_speakers(base_url, session), args.style_id)
    print(
        f"使用音声: {voice.speaker_name} / {voice.style_name} "
        f"(style_id={voice.style_id})"
    )

    audio_dir = args.output_dir / "slides"
    work_dir = args.output_dir / ".work"
    audio_dir.mkdir(parents=True, exist_ok=True)
    work_dir.mkdir(parents=True, exist_ok=True)

    wav_files: list[Path] = []
    manifest: list[dict[str, Any]] = []

    for index, (slide_number, text) in enumerate(notes, start=1):
        wav_path = work_dir / f"slide-{slide_number:03d}.wav"
        destination = audio_dir / f"slide-{slide_number:03d}.{args.format}"
        print(f"[{index}/{len(notes)}] スライド {slide_number} を音声化")
        query = make_audio_query(base_url, text, voice.style_id, args, session)
        synthesize(base_url, query, voice.style_id, wav_path, session)
        convert_audio(wav_path, destination)
        wav_files.append(wav_path)
        manifest.append(
            {
                "slide": slide_number,
                "notes_file": f"notes/slide-{slide_number:03d}.txt",
                "audio_file": f"slides/{destination.name}",
                "characters": len(text),
            }
        )

    manifest_path = args.output_dir / "manifest.json"
    manifest_path.write_text(
        json.dumps(
            {
                "source": args.pptx.name,
                "voice": {
                    "style_id": voice.style_id,
                    "speaker": voice.speaker_name,
                    "style": voice.style_name,
                },
                "slides": manifest,
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )

    if args.combine and wav_files:
        destination = args.output_dir / f"{args.combined_name}.{args.format}"
        combine_wavs(wav_files, args.slide_gap, destination, work_dir)
        print(f"結合音声を出力しました: {destination}")

    print(f"完了: {args.output_dir}")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except (ValueError, RuntimeError, requests.RequestException, subprocess.CalledProcessError) as exc:
        print(f"エラー: {exc}", file=sys.stderr)
        raise SystemExit(1)
